"""CWI Presentation Agent — generates a .pptx from any CWI agent JSON output."""

from __future__ import annotations

import json
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from rich.console import Console

from orchestrator.base_agent import BaseAgent
from orchestrator.settings import settings

console = Console(legacy_windows=False)

# ── Palette ───────────────────────────────────────────────────────────────────
_DARK  = RGBColor(0x1E, 0x1E, 0x2E)   # slide background
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_ACCENT = RGBColor(0x4A, 0x90, 0xD9)  # blue accent
_MUTED = RGBColor(0xAA, 0xAA, 0xBB)


class PresentationAgent(BaseAgent):
    name = "Presentation Agent"
    role = "Apresentador Executivo"
    model: str = settings.exec_report_model
    prompt_file = "agents/prompts_cwi/presentation.md"

    def __init__(self) -> None:
        self.model = settings.exec_report_model
        super().__init__()

    def run(self, input_file: str = "", input_data: dict[str, Any] | None = None) -> Path:
        """Generate a .pptx from a CWI agent JSON output file or dict."""
        console.rule("[bold]Presentation Agent")

        # 1. Load input
        if input_data:
            source = json.dumps(input_data, ensure_ascii=False, indent=2)
            source_label = "dict"
        elif input_file:
            source = Path(input_file).read_text(encoding="utf-8")
            source_label = Path(input_file).stem
        else:
            source, source_label = self._load_latest_output()

        # 2. Ask Claude to structure the slides
        user_message = f"""Analise o JSON abaixo e estruture uma apresentacao executiva em slides.

DADOS:
{source}

Retorne apenas o JSON de slides, sem texto adicional."""

        response_text = self._run(user_message, max_tokens=4096)
        slide_plan = self._parse_json_output(response_text)

        # 3. Build the .pptx
        pptx_path = self._build_pptx(slide_plan, source_label)

        console.print(f"[green]Apresentacao gerada:[/] {pptx_path}")
        return pptx_path

    # ── PPTX builder ──────────────────────────────────────────────────────────

    def _build_pptx(self, plan: dict[str, Any], label: str) -> Path:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # completely blank

        for slide_def in plan.get("slides", []):
            tipo = slide_def.get("tipo", "conteudo")
            slide = prs.slides.add_slide(blank_layout)
            self._fill_background(slide, _DARK)

            if tipo == "capa":
                self._render_capa(slide, slide_def, plan)
            elif tipo == "titulo_secao":
                self._render_secao(slide, slide_def)
            elif tipo == "encerramento":
                self._render_conteudo(slide, slide_def, accent=True)
            else:
                self._render_conteudo(slide, slide_def)

        today = date.today().strftime("%Y_%m_%d")
        out_path = Path(settings.output_dir) / "cwi" / f"presentation_{label}_{today}.pptx"
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        return out_path

    def _fill_background(self, slide: Any, color: RGBColor) -> None:
        from pptx.enum.dml import MSO_THEME_COLOR  # noqa: PLC0415
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text(
        self,
        slide: Any,
        text: str,
        left: float, top: float, width: float, height: float,
        font_size: int = 18,
        bold: bool = False,
        color: RGBColor = _WHITE,
        wrap: bool = True,
    ) -> Any:
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    def _render_capa(self, slide: Any, sd: dict, plan: dict) -> None:
        # accent bar
        from pptx.util import Emu  # noqa: PLC0415
        shape = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.33), Inches(0.06))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _ACCENT
        shape.line.fill.background()

        self._add_text(slide, sd.get("titulo", plan.get("titulo_apresentacao", "")),
                       0.6, 1.5, 12.0, 1.4, font_size=36, bold=True)
        self._add_text(slide, sd.get("subtitulo", plan.get("subtitulo", "")),
                       0.6, 3.4, 12.0, 0.8, font_size=20, color=_MUTED)
        self._add_text(slide, "CWI Software", 0.6, 6.6, 4.0, 0.6, font_size=14, color=_MUTED)

    def _render_secao(self, slide: Any, sd: dict) -> None:
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _ACCENT
        shape.line.fill.background()

        self._add_text(slide, sd.get("titulo", ""),
                       0.5, 2.8, 12.0, 1.2, font_size=32, bold=True)

    def _render_conteudo(self, slide: Any, sd: dict, accent: bool = False) -> None:
        title_color = _ACCENT if accent else _WHITE
        self._add_text(slide, sd.get("titulo", ""),
                       0.5, 0.3, 12.3, 0.8, font_size=24, bold=True, color=title_color)

        # divider line
        shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.03))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _ACCENT
        shape.line.fill.background()

        bullets = sd.get("bullets", [])
        top = 1.3
        for bullet in bullets[:6]:
            self._add_text(slide, f"• {bullet}", 0.7, top, 11.8, 0.7, font_size=16)
            top += 0.72

        notas = sd.get("notas", "")
        if notas:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notas

    # ── Fallback loader ───────────────────────────────────────────────────────

    @staticmethod
    def _load_latest_output() -> tuple[str, str]:
        out = Path("outputs/cwi/")
        candidates = []
        for prefix in ("executive_report_", "pmo_", "agile_coach_", "product_"):
            files = sorted(out.glob(f"{prefix}*.json"), key=lambda p: p.stat().st_mtime, reverse=True)
            if files:
                candidates.append(files[0])
                break
        if not candidates:
            raise FileNotFoundError("Nenhum output CWI encontrado. Rode um agente antes.")
        f = candidates[0]
        return f.read_text(encoding="utf-8"), f.stem
